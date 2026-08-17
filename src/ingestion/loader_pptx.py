"""Extração de texto de PPTX: conteúdo de cada slide + notas do apresentador."""
from pathlib import Path

from pptx import Presentation


def _slide_title(slide) -> str:
    if slide.shapes.title and slide.shapes.title.text:
        return slide.shapes.title.text.strip()
    return "(sem título)"


def _slide_body_text(slide, title: str) -> str:
    parts = []
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        text = shape.text_frame.text.strip()
        if not text or text == title:
            continue
        parts.append(text)
    return "\n".join(parts)


def load(file_path: Path) -> list[dict]:
    prs = Presentation(str(file_path))
    units: list[dict] = []

    for i, slide in enumerate(prs.slides, start=1):
        title = _slide_title(slide)
        body = _slide_body_text(slide, title)

        notes = ""
        if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
            notes = slide.notes_slide.notes_text_frame.text.strip()

        text = f"Slide {i} - Título: {title}\nConteúdo: {body}"
        if notes:
            text += f"\nNotas do apresentador: {notes}"

        units.append({"text": text, "location": f"slide {i}", "section": title})

    return units
