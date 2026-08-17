"""Gera os 18 documentos fictícios da Clínica Vida Plena em data/raw/<categoria>/.

Conteúdo 100% fictício, criado para fins de teste do pipeline RAG.
Rodar com o venv do projeto ativado: python scripts/generate_documents.py
"""
import csv
import json
from pathlib import Path

from docx import Document
from docx.shared import Pt
from openpyxl import Workbook
from pptx import Presentation
from pptx.util import Inches, Pt as PptPt
from reportlab.lib.pagesizes import A4
from reportlab.lib.styles import getSampleStyleSheet
from reportlab.platypus import Paragraph, SimpleDocTemplate, Spacer

import sys
sys.path.insert(0, str(Path(__file__).resolve().parent.parent))
from config import RAW_DOCS_DIR  # noqa: E402


def out(category: str, filename: str) -> Path:
    path = RAW_DOCS_DIR / category / filename
    path.parent.mkdir(parents=True, exist_ok=True)
    return path


# --------------------------------------------------------------------------- PDF
def make_pdf(path: Path, title: str, sections: list[tuple[str, str]]) -> None:
    styles = getSampleStyleSheet()
    doc = SimpleDocTemplate(str(path), pagesize=A4)
    story = [Paragraph(title, styles["Title"]), Spacer(1, 16)]
    for heading, body in sections:
        story.append(Paragraph(heading, styles["Heading2"]))
        story.append(Spacer(1, 6))
        for para in body.strip().split("\n\n"):
            story.append(Paragraph(para.strip().replace("\n", "<br/>"), styles["BodyText"]))
            story.append(Spacer(1, 8))
        story.append(Spacer(1, 10))
    doc.build(story)


def gen_manual_colaborador():
    sections = [
        ("1. Boas-vindas",
         "Seja bem-vindo(a) à Clínica Vida Plena! Este manual reúne as principais políticas de "
         "Recursos Humanos que todo colaborador deve conhecer, incluindo jornada de trabalho, "
         "período de experiência, benefícios e canais de suporte interno."),
        ("2. Jornada de trabalho",
         "A jornada padrão é de 44 horas semanais, distribuídas de segunda a sábado conforme "
         "escala da unidade. Colaboradores da recepção e do atendimento seguem escala de plantão "
         "definida pela Coordenação de Operações, com folgas compensatórias registradas em banco "
         "de horas.\n\n"
         "Atrasos superiores a 15 minutos devem ser justificados junto à liderança direta e "
         "registrados no sistema de ponto eletrônico no mesmo dia."),
        ("3. Período de experiência",
         "Todo colaborador contratado por prazo indeterminado passa por um período de experiência "
         "de 90 dias, dividido em dois ciclos de 45 dias. Ao final de cada ciclo, o gestor direto "
         "realiza uma avaliação formal registrada no sistema de RH."),
        ("4. Benefícios",
         "Os benefícios oferecidos incluem plano de saúde, vale-refeição ou vale-alimentação, "
         "vale-transporte e auxílio-creche para colaboradores com filhos até 6 anos. Detalhes "
         "completos de elegibilidade por cargo estão na tabela de benefícios (planilha anexa "
         "mantida pelo RH)."),
        ("5. Canais de suporte",
         "Dúvidas sobre folha de pagamento, benefícios ou políticas de RH devem ser direcionadas "
         "ao e-mail rh@clinicavidaplena.com.br ou ao ramal 4101."),
    ]
    make_pdf(out("rh", "manual_colaborador.pdf"), "Manual do Colaborador — Clínica Vida Plena", sections)


def gen_politica_reembolso_despesas():
    sections = [
        ("1. Objetivo",
         "Esta política define as regras para solicitação de reembolso de despesas incorridas por "
         "colaboradores no exercício de suas funções, como transporte, alimentação em viagens a "
         "trabalho e materiais de escritório de uso emergencial."),
        ("2. Limites de valor",
         "Despesas de alimentação em viagem têm limite diário de R$ 80,00. Despesas de transporte "
         "por aplicativo têm limite de R$ 60,00 por trajeto, exceto em deslocamentos entre unidades "
         "autorizados previamente pela liderança."),
        ("3. Prazo para solicitação",
         "O colaborador tem até 15 dias corridos após a despesa para enviar a solicitação de "
         "reembolso, acompanhada de nota fiscal ou recibo legível, através do sistema financeiro "
         "interno.\n\n"
         "Solicitações fora do prazo são analisadas caso a caso pelo Financeiro e podem ser "
         "recusadas."),
        ("4. Prazo de pagamento",
         "Reembolsos aprovados são pagos em até 10 dias úteis, na mesma data do fechamento da "
         "folha de pagamento mensal mais próxima."),
        ("5. Despesas não reembolsáveis",
         "Não são reembolsáveis despesas de caráter pessoal, multas de trânsito, bebidas "
         "alcoólicas e despesas sem comprovante fiscal válido."),
    ]
    make_pdf(out("financeiro", "politica_reembolso_despesas.pdf"),
             "Política de Reembolso de Despesas — Clínica Vida Plena", sections)


def gen_manual_atendimento_paciente():
    sections = [
        ("1. Fluxo de recepção",
         "Ao chegar à unidade, o paciente deve ser recebido pela recepção em até 5 minutos, com "
         "confirmação de dados cadastrais e do convênio ou forma de pagamento particular. "
         "Pacientes com consulta agendada têm prioridade de check-in em relação a encaixes."),
        ("2. Agendamento",
         "Consultas podem ser agendadas por telefone, WhatsApp institucional ou pelo portal do "
         "paciente. O intervalo padrão entre consultas é de 30 minutos, podendo variar conforme a "
         "especialidade médica."),
        ("3. Encaixes",
         "Encaixes de urgência são autorizados pela coordenação da unidade e não devem ultrapassar "
         "2 por período de atendimento, para não comprometer o tempo médio de espera dos demais "
         "pacientes."),
        ("4. Triagem inicial",
         "A triagem inicial é realizada pela equipe de enfermagem, com aferição de sinais vitais "
         "básicos e registro no prontuário eletrônico antes do atendimento médico."),
        ("5. Encerramento do atendimento",
         "Ao final da consulta, a recepção confirma o próximo agendamento (se houver) e orienta o "
         "paciente sobre emissão de nota fiscal e reembolso junto ao convênio, quando aplicável."),
    ]
    make_pdf(out("operacional", "manual_atendimento_paciente.pdf"),
             "Manual de Atendimento ao Paciente — Clínica Vida Plena", sections)


def gen_politica_privacidade_lgpd():
    sections = [
        ("1. Objetivo",
         "Esta política estabelece as diretrizes para tratamento de dados pessoais e dados "
         "sensíveis de saúde de pacientes, colaboradores e parceiros da Clínica Vida Plena, em "
         "conformidade com a Lei Geral de Proteção de Dados (Lei 13.709/2018)."),
        ("2. Dados coletados",
         "São coletados dados cadastrais (nome, CPF, endereço, contato), dados de convênio e dados "
         "sensíveis de saúde (histórico clínico, exames, prescrições), sempre vinculados a uma "
         "finalidade específica de atendimento."),
        ("3. Base legal e finalidade",
         "O tratamento de dados de saúde se baseia na tutela da saúde, conforme art. 11, II, 'f' "
         "da LGPD, sendo utilizados exclusivamente para prestação do serviço médico, faturamento "
         "junto a convênios e cumprimento de obrigações legais."),
        ("4. Compartilhamento",
         "Dados podem ser compartilhados com operadoras de convênio médico e laboratórios "
         "parceiros, mediante contrato que garanta o mesmo nível de proteção exigido pela clínica. "
         "Não há venda ou compartilhamento de dados para fins de marketing de terceiros."),
        ("5. Direitos do titular",
         "O paciente ou colaborador pode solicitar acesso, correção, portabilidade ou eliminação "
         "de seus dados através do canal compliance@clinicavidaplena.com.br, endereçado ao "
         "Encarregado de Proteção de Dados (DPO)."),
        ("6. Retenção",
         "Prontuários médicos são mantidos por no mínimo 20 anos, conforme resolução do Conselho "
         "Federal de Medicina, mesmo após o encerramento do vínculo do paciente com a clínica."),
    ]
    make_pdf(out("legal_compliance", "politica_privacidade_lgpd.pdf"),
             "Política de Privacidade e Proteção de Dados (LGPD) — Clínica Vida Plena", sections)


def gen_protocolo_biosseguranca():
    sections = [
        ("1. Objetivo",
         "Este protocolo define os procedimentos obrigatórios de biossegurança e controle de "
         "infecção nos consultórios e salas de procedimento da Clínica Vida Plena."),
        ("2. Higienização das mãos",
         "A higienização das mãos é obrigatória antes e após contato com cada paciente, seguindo "
         "os 5 momentos definidos pela Organização Mundial da Saúde. Álcool em gel 70% deve estar "
         "disponível em todas as salas."),
        ("3. Equipamentos de proteção individual (EPI)",
         "O uso de máscara cirúrgica é obrigatório durante todo o atendimento clínico. Luvas "
         "descartáveis devem ser trocadas a cada paciente e descartadas em lixo hospitalar "
         "identificado (saco branco leitoso)."),
        ("4. Esterilização de instrumentais",
         "Todo instrumental reutilizável passa por lavagem, desinfecção e esterilização em "
         "autoclave antes de novo uso, conforme checklist de esterilização (planilha anexa). "
         "Autoclaves são validadas mensalmente por empresa terceirizada."),
        ("5. Descarte de resíduos",
         "Resíduos perfurocortantes são descartados em coletores rígidos identificados, "
         "substituídos ao atingir 2/3 da capacidade, nunca por preenchimento total."),
    ]
    make_pdf(out("qualidade_biosseguranca", "protocolo_biosseguranca_controle_infeccao.pdf"),
             "Protocolo de Biossegurança e Controle de Infecção — Clínica Vida Plena", sections)


# --------------------------------------------------------------------------- DOCX
def add_heading_body(doc: Document, heading: str, body: str, level: int = 1):
    doc.add_heading(heading, level=level)
    for para in body.strip().split("\n\n"):
        doc.add_paragraph(para.strip())


def gen_politica_ferias_banco_horas():
    doc = Document()
    doc.add_heading("Política de Férias e Banco de Horas — Clínica Vida Plena", level=0)
    add_heading_body(doc, "1. Solicitação de férias",
        "As férias devem ser solicitadas com no mínimo 30 dias de antecedência através do sistema "
        "de RH, respeitando o período aquisitivo de 12 meses de trabalho.\n\n"
        "É vedado o fracionamento em mais de 3 períodos, sendo que um deles não pode ser inferior "
        "a 14 dias corridos, conforme legislação trabalhista vigente.")
    add_heading_body(doc, "2. Aprovação",
        "A liderança direta tem até 5 dias úteis para aprovar ou solicitar ajuste na data "
        "proposta, considerando a necessidade de cobertura mínima de plantão na unidade.")
    add_heading_body(doc, "3. Banco de horas",
        "Colaboradores da recepção e enfermagem que atuam em escala podem compensar horas extras "
        "via banco de horas, com saldo máximo acumulado de 40 horas.\n\n"
        "O saldo do banco de horas deve ser compensado em até 6 meses, sob pena de pagamento como "
        "hora extra no fechamento da folha seguinte.")
    add_heading_body(doc, "4. Plantões",
        "Plantões de fim de semana e feriado são remunerados com adicional de 50% sobre a hora "
        "normal, ou compensados em banco de horas mediante acordo individual registrado com o RH.")
    doc.save(out("rh", "politica_ferias_banco_horas.docx"))


def gen_procedimento_pagamento_fornecedores():
    doc = Document()
    doc.add_heading("Procedimento de Pagamento a Fornecedores — Clínica Vida Plena", level=0)
    add_heading_body(doc, "1. Cadastro de fornecedor",
        "Todo novo fornecedor deve ser cadastrado no sistema financeiro com CNPJ, dados bancários "
        "e contrato ou proposta comercial assinada, antes da emissão da primeira nota fiscal.")
    add_heading_body(doc, "2. Fluxo de aprovação",
        "Notas fiscais de até R$ 2.000,00 são aprovadas diretamente pelo gestor do centro de "
        "custo. Valores acima disso exigem aprovação adicional da Diretoria Executiva.\n\n"
        "O prazo máximo entre o recebimento da nota fiscal e o início da análise financeira é de "
        "3 dias úteis.")
    add_heading_body(doc, "3. Condições de pagamento",
        "O pagamento padrão a fornecedores ocorre em 28 dias após o recebimento e aprovação da "
        "nota fiscal, com remessas processadas toda quinta-feira.")
    add_heading_body(doc, "4. Centros de custo",
        "Cada despesa deve ser vinculada a um centro de custo válido (ver planilha de centros de "
        "custo mantida pelo Financeiro), para permitir o rateio correto nos relatórios gerenciais.")
    doc.save(out("financeiro", "procedimento_pagamento_fornecedores.docx"))


def gen_protocolo_agendamento_cancelamento():
    doc = Document()
    doc.add_heading("Protocolo de Agendamento e Cancelamento de Consultas — Clínica Vida Plena", level=0)
    add_heading_body(doc, "1. Confirmação de consulta",
        "Todas as consultas agendadas recebem confirmação automática por SMS ou WhatsApp 48 horas "
        "antes do horário marcado. A recepção deve registrar a confirmação (ou ausência de "
        "resposta) no sistema.")
    add_heading_body(doc, "2. Cancelamento pelo paciente",
        "O paciente pode cancelar ou remarcar sua consulta sem custo até 24 horas antes do "
        "horário agendado, pelo telefone da unidade ou pelo portal do paciente.\n\n"
        "Cancelamentos com menos de 24 horas de antecedência, ou faltas sem aviso, podem gerar "
        "cobrança de taxa de no-show para consultas particulares, conforme tabela vigente.")
    add_heading_body(doc, "3. Cancelamento pela clínica",
        "Em caso de indisponibilidade do médico, a recepção deve contatar o paciente com no "
        "mínimo 12 horas de antecedência, oferecendo reagendamento prioritário na mesma semana.")
    add_heading_body(doc, "4. Lista de espera",
        "Pacientes que aceitam entrar na lista de espera são chamados automaticamente pelo "
        "sistema em caso de cancelamento de outro paciente, respeitando a ordem de solicitação.")
    doc.save(out("operacional", "protocolo_agendamento_cancelamento.docx"))


def gen_termo_consentimento():
    doc = Document()
    doc.add_heading("Termo de Consentimento para Tratamento de Dados do Paciente — Clínica Vida Plena", level=0)
    add_heading_body(doc, "1. Finalidade do termo",
        "Este termo formaliza o consentimento do paciente para a coleta e o tratamento de seus "
        "dados pessoais e dados sensíveis de saúde, necessários à prestação do atendimento médico "
        "pela Clínica Vida Plena.")
    add_heading_body(doc, "2. Dados tratados",
        "Serão tratados dados cadastrais, histórico clínico, resultados de exames, prescrições e "
        "dados de contato, exclusivamente para fins assistenciais, administrativos e de "
        "faturamento junto a convênios.")
    add_heading_body(doc, "3. Compartilhamento com terceiros",
        "Os dados poderão ser compartilhados com laboratórios parceiros e operadoras de convênio, "
        "estritamente para viabilizar o atendimento solicitado pelo próprio paciente.")
    add_heading_body(doc, "4. Direitos garantidos",
        "O paciente pode, a qualquer momento, solicitar acesso, correção ou exclusão de seus "
        "dados (exceto quando houver obrigação legal de retenção do prontuário), através do canal "
        "compliance@clinicavidaplena.com.br.")
    add_heading_body(doc, "5. Assinatura",
        "Declaro estar ciente e de acordo com os termos acima descritos para o tratamento dos "
        "meus dados pessoais pela Clínica Vida Plena.")
    doc.save(out("legal_compliance", "termo_consentimento_dados_paciente.docx"))


# --------------------------------------------------------------------------- XLSX
def gen_tabela_beneficios_por_cargo():
    wb = Workbook()
    ws = wb.active
    ws.title = "Benefícios"
    ws.append(["Cargo", "Plano de Saúde", "Vale-Refeição (R$/dia)", "Vale-Transporte", "PLR"])
    rows = [
        ["Recepcionista", "Básico", 35, "Sim", "Sim, conforme meta da unidade"],
        ["Técnico de Enfermagem", "Básico", 35, "Sim", "Sim, conforme meta da unidade"],
        ["Enfermeiro(a)", "Intermediário", 40, "Sim", "Sim, conforme meta da unidade"],
        ["Médico(a) Plantonista", "Intermediário", 45, "Não aplicável", "Sim, conforme produtividade"],
        ["Coordenador(a) de Unidade", "Superior", 45, "Sim", "Sim, conforme meta da unidade"],
        ["Analista Financeiro", "Intermediário", 40, "Sim", "Sim, conforme resultado da empresa"],
        ["Gerente / Diretoria", "Superior", 50, "Não aplicável", "Sim, conforme resultado da empresa"],
    ]
    for row in rows:
        ws.append(row)
    wb.save(out("rh", "tabela_beneficios_por_cargo.xlsx"))


def gen_checklist_esterilizacao():
    wb = Workbook()
    ws = wb.active
    ws.title = "Checklist Autoclave"
    ws.append(["Item verificado", "Frequência", "Responsável", "Critério de conformidade"])
    rows = [
        ["Teste biológico da autoclave", "Semanal", "Enfermagem", "Resultado negativo para esporos"],
        ["Teste de Bowie-Dick", "Diário (1ª carga do dia)", "Enfermagem", "Viragem uniforme do indicador"],
        ["Validação técnica externa", "Mensal", "Empresa terceirizada", "Laudo de conformidade emitido"],
        ["Registro de ciclo (tempo/temperatura)", "A cada ciclo", "Enfermagem", "Dentro dos parâmetros do fabricante"],
        ["Integridade das embalagens esterilizadas", "A cada ciclo", "Enfermagem", "Sem furos, umidade ou violação"],
        ["Limpeza interna da câmara", "Semanal", "Enfermagem", "Sem resíduos visíveis"],
    ]
    for row in rows:
        ws.append(row)
    wb.save(out("qualidade_biosseguranca", "checklist_esterilizacao_autoclave.xlsx"))


# --------------------------------------------------------------------------- CSV
def gen_centros_de_custo():
    path = out("financeiro", "centros_de_custo.csv")
    with open(path, "w", newline="", encoding="utf-8") as f:
        writer = csv.writer(f)
        writer.writerow(["codigo_centro_custo", "nome", "area_responsavel", "responsavel"])
        writer.writerows([
            ["CC-100", "Diretoria Executiva", "Estratégico", "Camila Duarte"],
            ["CC-200", "Recursos Humanos", "RH", "Fernanda Lopes"],
            ["CC-300", "Financeiro e Contábil", "Financeiro", "Marcelo Andrade"],
            ["CC-400", "Atendimento e Recepção", "Operacional", "Juliana Ramos"],
            ["CC-410", "Enfermagem e Procedimentos", "Operacional", "Juliana Ramos"],
            ["CC-500", "Jurídico e Compliance", "Legal e Compliance", "Rodrigo Nascimento"],
            ["CC-600", "Qualidade e Biossegurança", "Qualidade", "Patrícia Menezes"],
            ["CC-700", "Comunicação Interna", "Comunicação Interna", "Bruno Teixeira"],
            ["CC-800", "Manutenção e Facilities", "Operacional", "Juliana Ramos"],
        ])


# --------------------------------------------------------------------------- JSON
def gen_faq_compliance_lgpd():
    faq = [
        {
            "pergunta": "Por quanto tempo a clínica guarda o prontuário do paciente?",
            "resposta": "O prontuário médico é mantido por no mínimo 20 anos, conforme resolução do "
                        "Conselho Federal de Medicina, mesmo após o fim do vínculo do paciente com a clínica."
        },
        {
            "pergunta": "Um colaborador pode acessar o prontuário de qualquer paciente?",
            "resposta": "Não. O acesso ao prontuário é restrito aos profissionais diretamente envolvidos "
                        "no atendimento daquele paciente, com registro de log de acesso auditável."
        },
        {
            "pergunta": "O que fazer se um paciente pedir a exclusão dos seus dados?",
            "resposta": "A solicitação deve ser encaminhada ao Encarregado de Proteção de Dados (DPO) via "
                        "compliance@clinicavidaplena.com.br. Dados de prontuário não podem ser excluídos "
                        "antes do prazo legal de retenção, mas dados de marketing e cadastro complementar "
                        "podem ser removidos mediante solicitação."
        },
        {
            "pergunta": "É permitido compartilhar dados de pacientes com laboratórios parceiros?",
            "resposta": "Sim, desde que exista contrato vigente que garanta o mesmo nível de proteção de "
                        "dados exigido pela clínica, e que o compartilhamento seja estritamente necessário "
                        "para o atendimento solicitado pelo paciente."
        },
        {
            "pergunta": "O que caracteriza um incidente de segurança de dados na clínica?",
            "resposta": "Qualquer acesso não autorizado, vazamento, perda ou alteração indevida de dados "
                        "pessoais ou de saúde. Incidentes devem ser reportados imediatamente ao Jurídico e "
                        "Compliance pelo ramal 4104."
        },
    ]
    path = out("legal_compliance", "faq_compliance_lgpd.json")
    with open(path, "w", encoding="utf-8") as f:
        json.dump({"documento": "FAQ de Compliance e LGPD — Clínica Vida Plena", "itens": faq}, f,
                  ensure_ascii=False, indent=2)


# --------------------------------------------------------------------------- HTML
def gen_comunicado_home_office():
    html = """<!DOCTYPE html>
<html lang="pt-br">
<head><meta charset="utf-8"><title>Comunicado Interno — Política de Home Office</title></head>
<body>
<nav>Intranet Clínica Vida Plena</nav>
<h1>Comunicado Interno: Nova Política de Home Office Administrativo</h1>
<p>Data de publicação: 03/02/2026 — Comunicação Interna</p>
<h2>O que muda</h2>
<p>A partir de março de 2026, colaboradores das áreas administrativas (Financeiro, RH, Jurídico e
Compliance, e Comunicação Interna) poderão trabalhar em regime de home office até 2 dias por
semana, mediante alinhamento prévio com a liderança direta.</p>
<h2>Quem não é elegível</h2>
<p>Colaboradores de atendimento direto ao paciente (recepção, enfermagem, corpo médico) não são
elegíveis a este regime, por exigência de presença física nas unidades.</p>
<h2>Como solicitar</h2>
<p>A solicitação deve ser feita através do sistema de RH, indicando os dias fixos da semana em
regime remoto, com aprovação do gestor direto em até 5 dias úteis.</p>
<h2>Dúvidas</h2>
<p>Em caso de dúvidas, entre em contato com o RH pelo e-mail rh@clinicavidaplena.com.br ou ramal
4101.</p>
<script>console.log('tracking pixel ignorado pela ingestão');</script>
</body>
</html>
"""
    out("comunicacao_interna", "comunicado_home_office.html").write_text(html, encoding="utf-8")


# --------------------------------------------------------------------------- Markdown
def gen_newsletter_mensal():
    md = """# Newsletter Mensal — Clínica Vida Plena

*Edição de Janeiro/2026 — Comunicação Interna*

## Destaques do mês

A Clínica Vida Plena atingiu a marca de 12.000 atendimentos realizados em 2025, um crescimento de
18% em relação ao ano anterior. Parabéns a todas as equipes envolvidas!

## Nova unidade

Foi confirmada a abertura da terceira unidade da clínica, na Zona Sul, com previsão de
inauguração para o segundo semestre de 2026. O processo seletivo para novos colaboradores começa
em abril.

## Treinamentos do mês

- Treinamento de biossegurança e controle de infecção (obrigatório para equipe assistencial)
- Treinamento de atendimento ao paciente para novos recepcionistas

## Aniversariantes

Parabéns aos colaboradores que celebram aniversário em janeiro! A lista completa está disponível
no mural de cada unidade.

## Canal de sugestões

Sugestões e feedbacks sobre esta newsletter podem ser enviados para
comunicacao@clinicavidaplena.com.br.
"""
    out("comunicacao_interna", "newsletter_mensal.md").write_text(md, encoding="utf-8")


def gen_missao_visao_valores():
    md = """# Missão, Visão e Valores — Clínica Vida Plena

## Missão

Oferecer atendimento médico humanizado e de excelência, com acesso facilitado e cuidado integral
à saúde dos nossos pacientes.

## Visão

Ser reconhecida, até 2028, como a rede de clínicas multiespecialidade mais confiável da região
metropolitana, referência em qualidade assistencial e experiência do paciente.

## Valores

### Ética e transparência

Atuamos com integridade em todas as relações — com pacientes, colaboradores, convênios e
fornecedores.

### Cuidado centrado no paciente

Cada decisão operacional leva em conta o impacto na experiência e na segurança do paciente.

### Melhoria contínua

Buscamos constantemente aprimorar processos, protocolos clínicos e indicadores de qualidade.

### Valorização das pessoas

Investimos no desenvolvimento e bem-estar de nossos colaboradores como base para um atendimento
de excelência.

## Metas estratégicas 2026

1. Reduzir o tempo médio de espera em 20%.
2. Ampliar em 30% o número de convênios aceitos.
3. Inaugurar a terceira unidade da clínica.
4. Certificar a unidade principal em programa de acreditação em qualidade.
"""
    out("estrategico", "missao_visao_valores.md").write_text(md, encoding="utf-8")


# --------------------------------------------------------------------------- PPTX
def add_slide(prs: Presentation, title: str, bullets: list[str], notes: str | None = None):
    layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(layout)
    slide.shapes.title.text = title
    body = slide.placeholders[1].text_frame
    body.text = bullets[0]
    for b in bullets[1:]:
        p = body.add_paragraph()
        p.text = b
    if notes:
        slide.notes_slide.notes_text_frame.text = notes
    return slide


def gen_treinamento_recepcao():
    prs = Presentation()
    add_slide(prs, "Treinamento: Fluxo de Atendimento na Recepção",
              ["Clínica Vida Plena", "Onboarding de novos colaboradores"],
              notes="Slide de abertura. Apresentador deve se apresentar e explicar a duração do "
                    "treinamento (aproximadamente 40 minutos).")
    add_slide(prs, "Check-in do paciente",
              ["Confirmar dados cadastrais e convênio",
               "Prazo máximo de 5 minutos para check-in",
               "Priorizar pacientes com consulta agendada sobre encaixes"],
              notes="Reforçar que o check-in lento é a principal causa de reclamação de pacientes, "
                    "segundo pesquisa de satisfação de 2025.")
    add_slide(prs, "Agendamento e encaixes",
              ["Intervalo padrão de 30 minutos entre consultas",
               "Máximo de 2 encaixes de urgência por período",
               "Encaixes extras exigem aprovação da coordenação"],
              notes="Explicar que encaixes não autorizados geram atraso em cadeia para os demais "
                    "pacientes do período.")
    add_slide(prs, "Cancelamentos e no-show",
              ["Cancelamento gratuito até 24h antes",
               "Menos de 24h ou falta pode gerar taxa de no-show",
               "Sempre oferecer reagendamento prioritário"],
              notes="Ver protocolo completo de agendamento e cancelamento no manual operacional.")
    add_slide(prs, "Dúvidas e suporte",
              ["Dúvidas operacionais: Coordenação de Operações, ramal 4103",
               "Dúvidas de sistema: suporte de TI interno"],
              notes="Encerrar o treinamento e abrir espaço para perguntas.")
    prs.save(out("operacional", "treinamento_fluxo_atendimento_recepcao.pptx"))


def gen_planejamento_estrategico():
    prs = Presentation()
    add_slide(prs, "Planejamento Estratégico 2026",
               ["Clínica Vida Plena", "Diretoria Executiva"],
               notes="Apresentação institucional aberta a todos os colaboradores.")
    add_slide(prs, "Resultados de 2025",
               ["12.000 atendimentos realizados (+18% vs. 2024)",
                "Duas unidades em operação",
                "NPS médio de 78 pontos"],
               notes="Destacar o crescimento como resultado do esforço coletivo de todas as áreas.")
    add_slide(prs, "Prioridades 2026",
               ["Reduzir tempo médio de espera em 20%",
                "Ampliar em 30% o número de convênios aceitos",
                "Inaugurar a terceira unidade (Zona Sul)",
                "Certificação em programa de acreditação em qualidade"],
               notes="Cada prioridade terá um responsável de área e indicadores acompanhados "
                     "trimestralmente.")
    add_slide(prs, "Investimento em pessoas",
               ["Novo programa de treinamento contínuo",
                "Revisão da tabela de benefícios por cargo",
                "Pesquisa de clima organizacional semestral"],
               notes="RH apresentará o cronograma detalhado em reunião específica no próximo mês.")
    prs.save(out("estrategico", "planejamento_estrategico_2026.pptx"))


def main():
    gen_manual_colaborador()
    gen_politica_reembolso_despesas()
    gen_manual_atendimento_paciente()
    gen_politica_privacidade_lgpd()
    gen_protocolo_biosseguranca()

    gen_politica_ferias_banco_horas()
    gen_procedimento_pagamento_fornecedores()
    gen_protocolo_agendamento_cancelamento()
    gen_termo_consentimento()

    gen_tabela_beneficios_por_cargo()
    gen_checklist_esterilizacao()

    gen_centros_de_custo()

    gen_faq_compliance_lgpd()

    gen_comunicado_home_office()

    gen_newsletter_mensal()
    gen_missao_visao_valores()

    gen_treinamento_recepcao()
    gen_planejamento_estrategico()

    print("18 documentos gerados em", RAW_DOCS_DIR)


if __name__ == "__main__":
    main()
